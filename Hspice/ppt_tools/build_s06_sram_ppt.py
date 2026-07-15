from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from PIL import Image


ROOT = Path(r"C:\Users\Tengfei\Desktop\Project_DTCO")
AUTO = ROOT / "AutoLDM"
SRC = ROOT / "S06_sram.pptx"
OUT = ROOT / "S06_sram_research_update_with_methods.pptx"
FIG = AUTO / "Hspice" / "sram_full_metrics" / "figures"
SUMMARY = AUTO / "Hspice" / "sram_full_metrics" / "data" / "summary_metrics.csv"

RED = RGBColor(166, 0, 0)
DARK = RGBColor(50, 50, 50)
GRAY = RGBColor(105, 105, 105)
BLUE = RGBColor(31, 119, 180)
ORANGE = RGBColor(214, 96, 39)


def require(path: Path) -> Path:
    if not path.exists():
        raise FileNotFoundError(path)
    return path


def add_textbox(slide, x, y, w, h, text, size=16, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_title(slide, title, page_no):
    add_textbox(slide, 0.25, 0.25, 10.6, 0.45, f"u {title}", size=22, bold=True, color=DARK)
    add_textbox(slide, 11.6, 6.55, 0.8, 0.2, f"< {page_no} >", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, bullets, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(5)
    return box


def add_metric_box(slide, x, y, w, h, label, value, note="", color=RED):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
    shape.line.color.rgb = RGBColor(220, 220, 220)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Microsoft YaHei"
    r2.font.size = Pt(10)
    r2.font.color.rgb = DARK
    if note:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = note
        r3.font.name = "Microsoft YaHei"
        r3.font.size = Pt(8)
        r3.font.color.rgb = GRAY
    return shape


def add_image_fit(slide, img_path: Path, x, y, w, h):
    img_path = require(img_path)
    with Image.open(img_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        width = w
        height = w / img_ratio
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * img_ratio
        left = x + (w - width) / 2
        top = y
    return slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_caption(slide, x, y, w, text):
    add_textbox(slide, x, y, w, 0.22, text, size=9, color=GRAY, align=PP_ALIGN.CENTER)


def add_definition_table(slide, x, y, w, h, headers, rows, col_widths=None, font_size=10):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = Inches(w * cw / total)
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(235, 235, 235)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Microsoft YaHei"
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = DARK
    for r_idx, row in enumerate(rows, 1):
        for c_idx, text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = text
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c_idx > 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = DARK
    return table_shape


def add_blank(prs, title, page_no):
    layout_idx = 4 if len(prs.slide_layouts) > 4 else len(prs.slide_layouts) - 1
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    add_title(slide, title, page_no)
    return slide


def update_slide5(slide):
    # Replace title with readable scientific summary title.
    for shape in slide.shapes:
        if hasattr(shape, "text") and "SRAM" in shape.text:
            shape.text = "SRAM指标总览：RC主要增加能耗，静态稳定性变化较小"
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Microsoft YaHei"
                    r.font.size = Pt(22)
                    r.font.bold = True
                    r.font.color.rgb = DARK
            break
    add_metric_box(slide, 0.75, 2.9, 1.55, 0.85, "HSNM No-RC", "268 mV", "RC: 267.5 mV", BLUE)
    add_metric_box(slide, 2.55, 2.9, 1.55, 0.85, "RSNM No-RC", "108 mV", "RC: 106 mV", BLUE)
    add_metric_box(slide, 4.35, 2.9, 1.55, 0.85, "Read delay", "+0.4 ps", "RC vs No-RC", ORANGE)
    add_metric_box(slide, 6.15, 2.9, 1.55, 0.85, "Write delay", "+1.2 ps", "RC vs No-RC", ORANGE)
    add_metric_box(slide, 7.95, 2.9, 1.55, 0.85, "Read energy", "1.44×", "0.029→0.043 fJ", RED)
    add_metric_box(slide, 9.75, 2.9, 1.55, 0.85, "Write energy", "2.92×", "0.039→0.114 fJ", RED)
    add_textbox(
        slide,
        0.75,
        4.1,
        10.8,
        0.8,
        "结论：当前RC候选网络对HSNM/RSNM影响很小；主要代价体现在读写延迟轻微上升，以及写能耗显著增加。",
        size=15,
        bold=False,
        color=DARK,
    )


def fix_out_of_bounds(prs):
    sw, sh = prs.slide_width, prs.slide_height
    margin = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.left < margin:
                shape.left = margin
            if shape.top < margin:
                shape.top = margin
            if shape.left + shape.width > sw:
                shape.width = max(1, sw - shape.left)
            if shape.top + shape.height > sh:
                shape.height = max(1, sh - shape.top)


def build():
    require(SRC)
    require(SUMMARY)
    for name in [
        "sram6t_norc_schematic.png",
        "sram6t_rc_candidate_network.png",
        "metric_extraction_annotation.png",
        "hold_snm_max_square_comparison.png",
        "ideal_rsnm_noise.png",
        "rc_candidate_rsnm_noise.png",
        "ideal_read_waveform.png",
        "rc_candidate_read_waveform.png",
        "ideal_read_bldiff.png",
        "rc_candidate_read_bldiff.png",
        "ideal_write_waveform.png",
        "rc_candidate_write_waveform.png",
        "ideal_write_trip.png",
        "rc_candidate_write_trip.png",
        "ideal_energy.png",
        "rc_candidate_energy.png",
    ]:
        require(FIG / name)

    shutil.copy2(SRC, OUT)
    prs = Presentation(str(OUT))
    fix_out_of_bounds(prs)

    update_slide5(prs.slides[4])

    s = add_blank(prs, "测试方法：从功能波形中提取稳定性、速度与能耗", 6)
    add_image_fit(s, FIG / "metric_extraction_annotation.png", 0.55, 0.85, 7.15, 5.2)
    add_bullets(
        s,
        7.95,
        1.0,
        3.9,
        4.9,
        [
            "静态稳定性：HSNM采用hold butterfly最大内嵌正方形；RSNM采用read-mode noise-source sweep。",
            "动态读：read disturb = 读窗口内低电平存储节点最大抬升；read delay = WL到BL-BLB=50 mV。",
            "动态写：write delay = WL到Q=0.5VDD；write-trip BL drop作为写能力proxy。",
            "能耗：E = ∫VDD·|I(VDD)|dt；所有指标来自同一HSPICE流程。",
        ],
        size=13,
    )

    s = add_blank(prs, "指标含义与测试条件：静态稳定性", 7)
    add_definition_table(
        s,
        0.55,
        1.0,
        11.1,
        4.55,
        ["指标", "物理含义", "测试条件", "提取/计算方法"],
        [
            [
                "HSNM",
                "保持状态下抵抗静态噪声的能力；数值越大，存储节点越不易被扰动翻转。",
                "Hold mode：WL=0；SRAM内部交叉耦合反相器做DC butterfly。",
                "在两个butterfly lobe内搜索最大内嵌正方形；HSNM取左右lobe中较小边长。",
            ],
            [
                "RSNM",
                "读操作中抵抗读破坏的静态噪声容限；直接反映读稳定性。",
                "Read mode：WL=VDD，BL=BLB=VDD；对内部反馈路径注入反向噪声VN。",
                "VN从0到0.5 V按1 mV扫描；首次出现Q/QB翻转的临界VN为读噪声容限。",
            ],
            [
                "Read disturb",
                "读操作期间，原本低电平存储节点被预充位线抬高的最大幅度。",
                "Read transient：初始Q=1/QB=0，BL/BLB预充到VDD，WL打开。",
                "read_disturb = max(V(QB))；该量不是RSNM，但可直观看读扰动强弱。",
            ],
            [
                "Read stability margin",
                "动态读扰动距离VDD/2翻转阈值的剩余空间；辅助理解瞬态安全裕量。",
                "使用同一read transient波形和read disturb提取结果。",
                "margin = VDD/2 - read_disturb；VDD=0.7 V。",
            ],
        ],
        col_widths=[1.15, 3.0, 3.25, 3.7],
        font_size=9,
    )
    add_textbox(s, 0.65, 5.85, 10.7, 0.5, "汇报口径：HSNM/RSNM为静态噪声容限；read disturb与read stability margin是瞬态读操作的辅助解释量。", size=13, bold=True, color=DARK)

    s = add_blank(prs, "指标含义与测试条件：速度、写能力与能耗", 8)
    add_definition_table(
        s,
        0.55,
        1.0,
        11.1,
        4.7,
        ["指标", "物理含义", "测试条件", "提取/计算方法"],
        [
            [
                "Read delay",
                "读出速度proxy；位线差分达到sense阈值所需时间。",
                "Read transient；WL上升打开access管，BL/BLB预充至VDD。",
                "read_delay = t(BL−BLB=50 mV) − t(WL=0.5VDD)。",
            ],
            [
                "Write delay",
                "写入速度；存储节点跨过中点电压的时间。",
                "Write-0 transient；BL拉低、BLB保持高，WL打开。",
                "write_delay = t(Q=0.5VDD) − t(WL=0.5VDD)。",
            ],
            [
                "Write-trip BL drop",
                "写能力proxy；写0时BL需从VDD下拉多少才能触发cell翻转。",
                "Write-trip DC sweep；逐步降低BL，观察Q/QB交叉点。",
                "BL_drop = VDD − VBL_at_crossing；不等同严格WSNM/WRM。",
            ],
            [
                "Read/write energy",
                "单cell读/写操作中从电源端消耗的能量。",
                "对应read/write transient窗口，统计VDD电源电流。",
                "E = ∫ VDD·abs(I(VDD_SRC)) dt，单位fJ。",
            ],
            [
                "Hold leakage",
                "保持状态下的静态漏电功耗。",
                "Hold mode稳态；WL=0，存储节点维持既定状态。",
                "P_leak = VDD·abs(I_VDD)，单位pW。",
            ],
        ],
        col_widths=[1.25, 3.0, 3.2, 3.65],
        font_size=8.7,
    )
    add_textbox(s, 0.65, 5.92, 10.7, 0.45, "注意：delay/energy为cell-level指标，不包含decoder、sense amplifier和write driver。", size=13, bold=True, color=DARK)

    s = add_blank(prs, "对比对象：No-RC单元与RC candidate网络", 9)
    add_image_fit(s, FIG / "sram6t_norc_schematic.png", 0.45, 0.9, 5.45, 4.55)
    add_image_fit(s, FIG / "sram6t_rc_candidate_network.png", 6.15, 0.9, 5.7, 4.55)
    add_caption(s, 0.7, 5.5, 4.9, "No-RC：理想6T SRAM单元，作为器件/电路基线")
    add_caption(s, 6.2, 5.5, 5.6, "RC candidate：加入候选寄生网络，用于评估互连代价")
    add_textbox(s, 0.75, 6.05, 10.8, 0.35, "固定条件：VDD=0.7 V，T=25 °C，L=16 nm，WPG/WPD/WPU=25 nm，CQ/CQB=1 fF，CBL/CBLB=10 fF。", size=12, color=GRAY)

    s = add_blank(prs, "静态稳定性：RC对HSNM/RSNM影响较小", 10)
    add_image_fit(s, FIG / "hold_snm_max_square_comparison.png", 0.35, 0.85, 5.8, 4.25)
    add_image_fit(s, FIG / "ideal_rsnm_noise.png", 6.35, 0.78, 2.7, 2.25)
    add_image_fit(s, FIG / "rc_candidate_rsnm_noise.png", 9.05, 0.78, 2.7, 2.25)
    add_metric_box(s, 6.45, 3.35, 1.65, 0.8, "HSNM", "−0.5 mV", "268→267.5", BLUE)
    add_metric_box(s, 8.25, 3.35, 1.65, 0.8, "RSNM", "−2 mV", "108→106", BLUE)
    add_metric_box(s, 10.05, 3.35, 1.65, 0.8, "Leakage", "+0.21%", "1383.9→1386.8 pW", GRAY)
    add_textbox(s, 6.45, 4.45, 5.25, 0.9, "解释：RC候选网络没有显著破坏双稳态；读噪声容限降低约2 mV，主要表现为轻微读稳定性退化。", size=13)

    s = add_blank(prs, "读操作：RC轻微增加读扰动和读延迟", 11)
    add_image_fit(s, FIG / "ideal_read_waveform.png", 0.45, 0.85, 5.45, 2.45)
    add_image_fit(s, FIG / "rc_candidate_read_waveform.png", 6.15, 0.85, 5.45, 2.45)
    add_image_fit(s, FIG / "ideal_read_bldiff.png", 0.45, 3.65, 5.45, 2.2)
    add_image_fit(s, FIG / "rc_candidate_read_bldiff.png", 6.15, 3.65, 5.45, 2.2)
    add_metric_box(s, 0.65, 6.05, 2.0, 0.55, "Read disturb", "+1.5 mV", "143.7→145.2", ORANGE)
    add_metric_box(s, 3.0, 6.05, 2.0, 0.55, "Read delay", "+0.4 ps", "15.2→15.6", ORANGE)
    add_textbox(s, 5.35, 6.08, 5.9, 0.45, "BL-BLB=50 mV作为读延迟判据；RC使位线差分建立略慢。", size=12, color=GRAY)

    s = add_blank(prs, "写操作：写延迟小幅增加，write-trip略有改善", 12)
    add_image_fit(s, FIG / "ideal_write_waveform.png", 0.45, 0.85, 5.45, 2.55)
    add_image_fit(s, FIG / "rc_candidate_write_waveform.png", 6.15, 0.85, 5.45, 2.55)
    add_image_fit(s, FIG / "ideal_write_trip.png", 0.45, 3.72, 5.45, 2.15)
    add_image_fit(s, FIG / "rc_candidate_write_trip.png", 6.15, 3.72, 5.45, 2.15)
    add_metric_box(s, 0.65, 6.05, 2.0, 0.55, "Write delay", "+1.2 ps", "37.4→38.6", ORANGE)
    add_metric_box(s, 3.0, 6.05, 2.0, 0.55, "BL drop", "−4 mV", "454→450", BLUE)
    add_textbox(s, 5.35, 6.08, 5.9, 0.45, "write-trip BL drop为写能力proxy，不等同于严格WSNM/WRM。", size=12, color=GRAY)

    s = add_blank(prs, "能耗结论：RC对写能耗最敏感", 13)
    add_image_fit(s, FIG / "ideal_energy.png", 0.55, 0.9, 5.25, 3.7)
    add_image_fit(s, FIG / "rc_candidate_energy.png", 6.05, 0.9, 5.25, 3.7)
    add_metric_box(s, 0.85, 4.95, 2.1, 0.8, "Read energy", "1.44×", "0.029→0.043 fJ", ORANGE)
    add_metric_box(s, 3.3, 4.95, 2.1, 0.8, "Write energy", "2.92×", "0.039→0.114 fJ", RED)
    add_metric_box(s, 5.75, 4.95, 2.1, 0.8, "RSNM", "−2 mV", "small penalty", BLUE)
    add_metric_box(s, 8.2, 4.95, 2.1, 0.8, "Delay", "+0.4/+1.2 ps", "read/write", ORANGE)
    add_textbox(s, 0.9, 6.05, 10.3, 0.45, "汇报结论：在当前候选RC网络下，稳定性基本保持；主要优化重点应放在降低写能耗和检查RC网络连接完整性。", size=13, bold=True, color=DARK)

    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    build()
